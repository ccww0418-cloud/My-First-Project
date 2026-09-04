#!/usr/bin/env python3
"""BookBot 발표 자료 생성 — 7장, 16:9

실행:
    /tmp/pptenv/bin/python scripts/make-deck.py
결과:
    BookBot-발표자료.pptx

모든 수치는 코드/배포에서 확인한 값입니다. 임의로 바꾸지 마세요.
"""

from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUT = "BookBot-발표자료.pptx"

# ── 색 ────────────────────────────────────────────────────────────
BG     = RGBColor(0xFA, 0xF7, 0xF0)   # 한지
INK    = RGBColor(0x1C, 0x19, 0x17)
MUTED  = RGBColor(0x78, 0x71, 0x6A)
RED    = RGBColor(0xA0, 0x3A, 0x2C)   # 강조
BLUE   = RGBColor(0x1F, 0x5C, 0x8B)   # AWS
GREEN  = RGBColor(0x3F, 0x63, 0x45)   # 외부 API
LINE   = RGBColor(0xC9, 0xC0, 0xB2)
BOXBG  = RGBColor(0xFF, 0xFF, 0xFF)
TINT   = RGBColor(0xF2, 0xEC, 0xE0)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

KO = "Apple SD Gothic Neo"

W, H = In(13.333), In(7.5)


# ── 헬퍼 ──────────────────────────────────────────────────────────
def deck():
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background(); bg.shadow.inherit = False
    return s


def txt(s, x, y, w, h, runs, size=14, color=INK, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6, line=1.25):
    """runs: 문자열 또는 [(텍스트, {size,color,bold}), ...] 리스트의 리스트(줄 단위)"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    lines = [runs] if isinstance(runs, str) else runs
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        p.line_spacing = line
        parts = [(ln, {})] if isinstance(ln, str) else ln
        for t, o in parts:
            r = p.add_run(); r.text = t
            f = r.font
            f.name = KO
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", bold)
            f.color.rgb = o.get("color", color)
    return tb


def box(s, x, y, w, h, fill=BOXBG, edge=LINE, wpt=0.75, shape=MSO_SHAPE.RECTANGLE):
    b = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        b.fill.background()
    else:
        b.fill.solid(); b.fill.fore_color.rgb = fill
    if edge is None:
        b.line.fill.background()
    else:
        b.line.color.rgb = edge; b.line.width = Pt(wpt)
    b.shadow.inherit = False
    b.text_frame.text = ""
    return b


def card(s, x, y, w, h, title, sub=None, body=(), accent=INK, fill=BOXBG, tsize=13):
    box(s, x, y, w, h, fill=fill)
    pad = In(0.13)
    lines = [[(title, {"size": tsize, "bold": True, "color": accent})]]
    if sub:
        lines.append([(sub, {"size": 9.5, "color": RED})])
    for b in body:
        lines.append([(b, {"size": 9, "color": MUTED})])
    txt(s, x + pad, y + In(0.10), w - pad * 2, h - In(0.16), lines, space=2, line=1.18)


def rule(s, x, y, w, color=LINE, wpt=0.75):
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Emu(int(Pt(wpt))))
    ln.fill.solid(); ln.fill.fore_color.rgb = color
    ln.line.fill.background(); ln.shadow.inherit = False
    return ln


def arrow(s, x, y, w, h, color=MUTED, shape=MSO_SHAPE.RIGHT_ARROW):
    a = s.shapes.add_shape(shape, x, y, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background(); a.shadow.inherit = False
    return a


def head(s, num, title, kicker=None):
    txt(s, In(0.62), In(0.42), In(1.0), In(0.3),
        [[(num, {"size": 11, "bold": True, "color": RED})]])
    txt(s, In(1.05), In(0.36), In(9.4), In(0.45),
        [[(title, {"size": 23, "bold": True})]])
    if kicker:
        txt(s, In(1.05), In(0.86), In(11.6), In(0.3),
            [[(kicker, {"size": 11.5, "color": MUTED})]])
    rule(s, In(0.62), In(1.22), In(12.1))


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text.strip()


def table(s, x, y, w, rows, colw, hdr=True, fs=10, rh=In(0.30), hfs=10):
    """rows: [[cell,...], ...]  colw: 비율 리스트"""
    n, m = len(rows), len(rows[0])
    tot = sum(colw)
    xs, acc = [], x
    for c in colw:
        xs.append(acc); acc += Emu(int(w * c / tot))
    ws = [(xs[i + 1] - xs[i]) if i + 1 < m else (x + w - xs[i]) for i in range(m)]

    for ri, row in enumerate(rows):
        yy = y + Emu(int(rh * ri))
        if hdr and ri == 0:
            bar = box(s, x, yy, w, rh, fill=TINT, edge=None)
        for ci, cell in enumerate(row):
            isH = hdr and ri == 0
            b, col, sz = (True, INK, hfs) if isH else (False, INK, fs)
            if isinstance(cell, tuple):
                cell, o = cell
                b = o.get("bold", b); col = o.get("color", col); sz = o.get("size", sz)
            txt(s, xs[ci] + In(0.08), yy + In(0.055), ws[ci] - In(0.10), rh,
                [[(str(cell), {"size": sz, "bold": b, "color": col})]],
                space=0, line=1.0)
        if not (hdr and ri == 0):
            rule(s, x, yy, w, color=LINE, wpt=0.5)
    rule(s, x, y + Emu(int(rh * n)), w, color=LINE, wpt=0.5)
    return y + Emu(int(rh * n))


def mono(s, x, y, w, h, text, size=9.5, fill=RGBColor(0xF4, 0xF1, 0xE8), color=INK):
    box(s, x, y, w, h, fill=fill, edge=LINE)
    tb = s.shapes.add_textbox(x + In(0.14), y + In(0.11), w - In(0.28), h - In(0.20))
    tf = tb.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(text.strip("\n").split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(0); p.line_spacing = 1.18
        r = p.add_run(); r.text = ln
        r.font.name = "Menlo"; r.font.size = Pt(size); r.font.color.rgb = color


# ══════════════════════════════════════════════════════════════════
prs = deck()

# ── 1. 표지 ───────────────────────────────────────────────────────
SHOT = "docs/site-shot.png"   # 여기에 사이트 캡처를 두면 표지에 자동으로 들어갑니다

s = blank(prs)
box(s, 0, 0, In(0.30), H, fill=RED, edge=None)

txt(s, In(1.05), In(1.05), In(6.2), In(0.4),
    [[("Un Livre Pour Vous", {"size": 12, "color": MUTED})]])
txt(s, In(1.05), In(1.44), In(6.2), In(1.0),
    [[("BookBot", {"size": 44, "bold": True})]])
txt(s, In(1.05), In(2.42), In(6.3), In(1.1),
    [[("기분과 상황으로 책을 찾아주는", {"size": 17})],
     [("AWS 서버리스 챗봇", {"size": 17})]],
    space=2)
rule(s, In(1.05), In(3.52), In(3.4), color=INK, wpt=1.25)
txt(s, In(1.05), In(3.72), In(6.3), In(0.8),
    [[("추천한 책이 실제로 존재하는지", {"size": 12, "color": MUTED})],
     [("외부 서점 · 도서관 API 6곳으로 검증합니다", {"size": 12, "color": MUTED})]],
    space=2)

stats = [("백엔드", "8,175줄"), ("도서 API", "6곳"), ("LLM 도구", "5종"),
         ("프론트엔드", "2,222줄"), ("자동 검증", "474건"), ("AWS 서비스", "10종")]
for i, (k, v) in enumerate(stats):
    cx = In(1.05) + Emu(int(In(1.78) * (i % 3)))
    cy = In(4.72) + Emu(int(In(0.74) * (i // 3)))
    txt(s, cx, cy, In(1.7), In(0.28), [[(k, {"size": 9, "color": MUTED})]], space=1)
    txt(s, cx, cy + In(0.24), In(1.7), In(0.4),
        [[(v, {"size": 16, "bold": True, "color": RED})]])

box(s, In(1.05), In(6.42), In(5.35), In(0.62), fill=TINT, edge=None)
txt(s, In(1.28), In(6.55), In(5.0), In(0.4),
    [[("d2cmff9bta4e7l.cloudfront.net", {"size": 12, "bold": True, "color": BLUE})],
     [("공개 접속 · 로그인 없음 · 쓰지 않으면 요금 0원", {"size": 9, "color": MUTED})]],
    space=2)

# 사이트 캡처 — 있으면 넣고, 없으면 자리와 방법을 표시
SX, SY, SW_, SH_ = In(6.95), In(1.05), In(5.78), In(5.99)
if os.path.exists(SHOT):
    from PIL import Image  # 선택 의존성
    iw, ih = Image.open(SHOT).size
    scale = min(SW_ / iw, SH_ / ih)
    w, h = Emu(int(iw * scale)), Emu(int(ih * scale))
    pic = s.shapes.add_picture(SHOT, SX + Emu(int((SW_ - w) / 2)),
                               SY + Emu(int((SH_ - h) / 2)), w, h)
    pic.line.color.rgb = LINE
    pic.line.width = Pt(0.75)
else:
    box(s, SX, SY, SW_, SH_, fill=RGBColor(0xF4, 0xF1, 0xE8), edge=LINE)
    txt(s, SX + In(0.5), SY + In(2.4), SW_ - In(1.0), In(1.3),
        [[("사이트 화면 캡처 자리", {"size": 13, "bold": True, "color": MUTED})],
         [("", {"size": 6})],
         [("d2cmff9bta4e7l.cloudfront.net 을 열고 ⌘⇧4 로 캡처해서", {"size": 10, "color": MUTED})],
         [("docs/site-shot.png 로 저장한 뒤 이 스크립트를 다시 실행하세요", {"size": 10, "color": MUTED})]],
        align=PP_ALIGN.CENTER, space=3)

notes(s, """
BookBot 은 대화로 책을 추천하는 서버리스 챗봇입니다.

핵심 한 줄은 "추천한 책이 실제로 존재하는지 외부에서 검증한다" 입니다.
일반 LLM 챗봇은 모델이 기억하는 것을 말하는데, 그러면 없는 책을 그럴듯하게
만들어 냅니다. 이 서비스는 그것을 구조로 막았고, 오늘 발표의 절반이 그 이야기입니다.

숫자 여섯 개만 기억해 주세요. 도서 API 6곳, LLM 도구 5종, 자동 검증 474건.
검증 474건은 네트워크 없이 몇 초 안에 도는 것만 센 값입니다.
""")

# ── 2. 무엇이 다른가 ──────────────────────────────────────────────
s = blank(prs)
head(s, "01", "일반 챗봇과 무엇이 다른가", "모델이 아는 것을 말하지 않고, 외부에서 확인된 것만 말합니다")

y = In(1.62)
box(s, In(0.62), y, In(5.85), In(1.55), fill=BOXBG)
txt(s, In(0.85), y + In(0.18), In(5.4), In(0.4),
    [[("일반 챗봇", {"size": 14, "bold": True, "color": MUTED})]])
mono(s, In(0.85), y + In(0.62), In(5.4), In(0.62),
     "입력 → 모델 → 출력", size=12, fill=BG, color=MUTED)
txt(s, In(0.85), y + In(1.24), In(5.4), In(0.3),
    [[("모델이 기억하는 것을 말합니다", {"size": 10, "color": MUTED})]])

box(s, In(6.88), y, In(5.85), In(1.55), fill=BOXBG, edge=RED, wpt=1.25)
txt(s, In(7.11), y + In(0.18), In(5.4), In(0.4),
    [[("BookBot", {"size": 14, "bold": True, "color": RED})]])
mono(s, In(7.11), y + In(0.62), In(5.4), In(0.62),
     "입력 → 검문 → 모델 ↔ 외부 DB → 검증 → 출력", size=11, fill=BG)
txt(s, In(7.11), y + In(1.24), In(5.4), In(0.3),
    [[("외부에서 확인된 것만 말합니다", {"size": 10, "color": RED, "bold": True})]])

y2 = In(3.48)
txt(s, In(0.62), y2, In(12.1), In(0.35),
    [[("역할을 둘로 쪼갠 것이 설계의 전부입니다", {"size": 14, "bold": True})]])

table(s, In(0.62), y2 + In(0.45), In(12.1),
      [["", "무엇을 하나", "어디서 오나"],
       [("LLM (Claude Sonnet 4.6)", {"bold": True}), "왜 이 책인지 설명하는 말만 씁니다", "모델"],
       [("도서 API 6곳", {"bold": True, "color": RED}),
        ("제목 · 저자 · 표지 · 평점 · ISBN · 구매 링크", {"color": RED}),
        ("외부 서점 · 도서관", {"color": RED})]],
      colw=[3.0, 5.6, 3.4], rh=In(0.42), fs=11.5, hfs=10.5)

box(s, In(0.62), In(5.42), In(12.1), In(1.28), fill=TINT, edge=None)
txt(s, In(0.95), In(5.62), In(11.4), In(0.95),
    [[("카드에 보이는 모든 값이 실제 API 응답입니다. 그래서 환각이 구조적으로 불가능합니다.", {"size": 13, "bold": True})],
     [("\"없는 책을 추천하지 마세요\" 는 프롬프트에 쓴 부탁입니다. 검증 계층은 보장입니다.", {"size": 11, "color": MUTED})]],
    space=6)

notes(s, """
이 슬라이드가 전체 발표의 뼈대입니다.

LLM 에게 책을 고르게 하면 존재하지 않는 책이 나옵니다. 제목이 그럴듯하고
저자도 실제 인물이라 사용자는 검증할 방법이 없습니다.

그래서 역할을 쪼갰습니다. 모델은 "왜 이 책이 당신에게 맞는지" 설명하는 말만
씁니다. 제목·저자·표지·평점·ISBN·구매 링크는 전부 외부 도서 API 에서 옵니다.
화면의 책 카드는 API 응답을 그대로 렌더링한 것입니다.

프롬프트에 "없는 책을 추천하지 마세요" 라고 쓰는 건 부탁입니다. 지켜질 때도
있고 안 지켜질 때도 있습니다. 검증 계층은 보장입니다. 이 차이가 중요합니다.
""")

# ── 3. 아키텍처 ───────────────────────────────────────────────────
s = blank(prs)
head(s, "02", "아키텍처", "서버를 직접 관리하지 않는 서버리스 구성 · 리전 us-east-1")

# AWS 경계
box(s, In(2.55), In(1.52), In(10.18), In(3.62), fill=None, edge=BLUE, wpt=1.0)
txt(s, In(2.72), In(1.62), In(4.0), In(0.28),
    [[("AWS 클라우드", {"size": 10, "bold": True, "color": BLUE})]])

# 사용자
card(s, In(0.62), In(2.62), In(1.62), In(0.95), "사용자", "웹 브라우저", ("HTTPS",))
arrow(s, In(2.28), In(2.98), In(0.22), In(0.20))

# CloudFront
card(s, In(2.75), In(2.15), In(2.05), In(1.85), "CloudFront",
     "유일한 공개 진입점",
     ("오리진 2개", "WAF 부착", "OAC SigV4", "TLS 1.2+"), accent=BLUE)

# 분기 라벨
txt(s, In(4.92), In(2.28), In(1.0), In(0.3), [[("/*", {"size": 9.5, "color": MUTED})]])
txt(s, In(4.92), In(3.36), In(1.2), In(0.3), [[("/api/*", {"size": 9.5, "color": RED})]])
arrow(s, In(4.84), In(2.52), In(0.20), In(0.18))
arrow(s, In(4.84), In(3.60), In(0.20), In(0.18), color=RED)

# S3
card(s, In(5.12), In(2.15), In(1.95), In(0.85), "S3", "정적 사이트",
     ("버킷 완전 비공개",), accent=BLUE)
# API GW
card(s, In(5.12), In(3.28), In(1.95), In(0.85), "API Gateway", "HTTP API",
     ("타임아웃 30초 고정",), accent=BLUE)
arrow(s, In(7.11), In(3.60), In(0.20), In(0.18), color=RED)

# Lambda
card(s, In(7.38), In(2.95), In(1.95), In(1.55), "Lambda", "애플리케이션 전체",
     ("Node 22 · arm64", "1024MB · 90초", "예약 동시성 10"), accent=BLUE)

# 오른쪽 3개
for i, (t, sb, bd) in enumerate([
        ("Bedrock", "Claude Sonnet 4.6", ("ConverseStream", "도구 사용")),
        ("DynamoDB", "테이블 1개 4용도", ("세션·캐시·기록", "레이트리밋")),
        ("SSM", "API 키 5개", ("SecureString",))]):
    card(s, In(9.62), In(1.90) + Emu(int(In(1.08) * i)), In(2.95), In(0.92),
         t, sb, bd, accent=BLUE)
    arrow(s, In(9.36), In(2.26) + Emu(int(In(1.08) * i)), In(0.20), In(0.18), color=MUTED)

# 외부 API 경계
box(s, In(0.62), In(5.38), In(12.1), In(1.42), fill=None, edge=GREEN, wpt=1.0)
txt(s, In(0.80), In(5.48), In(7.0), In(0.28),
    [[("외부 인터넷 — 서점 · 도서관 · 무료전자책 6곳 (병렬 호출, 하나가 죽어도 나머지로 답변)",
       {"size": 10, "bold": True, "color": GREEN})]])
srcs = [("알라딘", "국내 서점"), ("국립중앙도서관", "국내 납본"), ("Google Books", "커버리지"),
        ("Open Library", "주제 분류"), ("Hardcover", "평점·무드"), ("Gutendex", "무료 전문")]
for i, (n, r) in enumerate(srcs):
    cx = In(0.85) + Emu(int(In(1.965) * i))
    card(s, cx, In(5.82), In(1.83), In(0.80), n, r, (), accent=INK, tsize=10.5)

arrow(s, In(8.26), In(4.58), In(0.18), In(0.72), color=GREEN, shape=MSO_SHAPE.DOWN_ARROW)

notes(s, """
왼쪽에서 오른쪽으로 읽으시면 됩니다.

CloudFront 가 유일한 공개 진입점입니다. 경로로 갈립니다 — 화면 파일은 S3 로,
/api/* 로 오는 요청만 API Gateway 를 거쳐 Lambda 로 갑니다. S3 버킷은 완전
비공개이고 CloudFront 만 읽을 수 있습니다.

Lambda 가 애플리케이션 전체입니다. Bedrock 으로 모델을 부르고, 도서 API 6곳을
병렬로 조회하고, DynamoDB 에 세션·캐시·기록·레이트리밋을 넣고, SSM 에서 API 키를
읽습니다.

아래 초록 박스가 오늘 강조할 부분입니다. 서점·도서관 6곳을 동시에 부릅니다.
Promise.allSettled 를 쓰기 때문에 한 곳이 죽어도 나머지로 답변합니다. 실제로
개발 중에 Gutendex 가 다운됐는데 Open Library 폴백이 동작했습니다.

서버가 없다는 말의 의미 — 항상 켜둔 컴퓨터가 없고 요청이 올 때만 실행됩니다.
쓰지 않으면 요금이 0원이고, 고정비는 WAF 월 7달러 하나뿐입니다.
""")

# ── 4. 요청 흐름 ──────────────────────────────────────────────────
s = blank(prs)
head(s, "03", "요청 한 건이 지나는 6단계", "\"조용히 위로가 되는 한국 소설 추천해주세요\" 를 예로")

steps = [
    ("1", "문지기", "오리진 비밀 헤더 · 2,000자 제한 · 세션ID UUID 검증 · IP 레이트리밋(분당 10 / 하루 150)", MUTED),
    ("2", "정책 검사", "정규식 27개 → 걸리면 모델 호출 0회 · 0.2초 종료   통과하면 의도 분류(BOOK / SERVICE / ATTACK)", MUTED),
    ("3", "도구 루프", "Bedrock 호출 → 도구 실행 → 결과 되돌려주기 (최대 4회) · 도구 5종", BLUE),
    ("4", "검증", "서점·도서관 6곳 조회 → ISBN13 중복 제거 → 제목·저자 대조 → 통과한 책만 남김", RED),
    ("5", "카드 선별", "답변이 언급한 책만 카드로 → 12장까지 채움 → 답변에 있는데 카드 없으면 역으로 재조회", RED),
    ("6", "출력", "SSE 전송 + DynamoDB 기록 저장(90일 TTL)", MUTED),
]
y = In(1.58)
for i, (n, t, d, c) in enumerate(steps):
    yy = y + Emu(int(In(0.72) * i))
    nb = box(s, In(0.62), yy, In(0.42), In(0.42),
             fill=(RED if c == RED else (BLUE if c == BLUE else TINT)), edge=None)
    txt(s, In(0.62), yy + In(0.09), In(0.42), In(0.3),
        [[(n, {"size": 12, "bold": True,
               "color": (WHITE if c in (RED, BLUE) else INK)})]], align=PP_ALIGN.CENTER)
    txt(s, In(1.22), yy + In(0.02), In(2.1), In(0.35), [[(t, {"size": 14, "bold": True})]])
    txt(s, In(3.32), yy + In(0.07), In(9.4), In(0.5), [[(d, {"size": 10.5, "color": MUTED})]])
    if i < len(steps) - 1:
        rule(s, In(0.62), yy + In(0.60), In(12.1))

box(s, In(0.62), In(6.02), In(5.9), In(0.78), fill=TINT, edge=None)
txt(s, In(0.88), In(6.18), In(5.5), In(0.5),
    [[("2단계에서 걸리면 0.2초에 끝납니다", {"size": 11.5, "bold": True})],
     [("정규식이라 네트워크 0회 · 모델 비용 0", {"size": 9.5, "color": MUTED})]], space=3)

box(s, In(6.83), In(6.02), In(5.9), In(0.78), fill=TINT, edge=None)
txt(s, In(7.09), In(6.18), In(5.5), In(0.5),
    [[("4·5단계가 이 서비스의 핵심입니다", {"size": 11.5, "bold": True, "color": RED})],
     [("모델이 고른 책을 외부 데이터로 걸러내는 구간", {"size": 9.5, "color": MUTED})]], space=3)

notes(s, """
1단계 문지기. CloudFront 가 주입한 비밀 헤더가 없으면 403 입니다. 함수 URL 을
직접 부르는 우회를 막습니다.

2단계 정책 검사가 2단입니다. 먼저 정규식 27개 — 미성년 보호 7, 프롬프트 인젝션 18,
개인정보 2. 여기서 걸리면 Bedrock 을 아예 부르지 않습니다. 실측 0.2초, 비용 0입니다.
통과하면 모델로 의도를 분류합니다. 주제를 보는 게 아니라 "무엇을 해달라는
요청인가" 만 봅니다.

3단계 도구 루프. 모델이 "이 도구를 이 인자로 불러줘" 를 구조화된 형태로
돌려주고, 결과를 되돌려주면 이어서 생각합니다. 코드 상한은 4회입니다.

4·5단계가 오늘 강조할 부분입니다. 다음 두 장에서 자세히 봅니다.
""")

# ── 5. 외부 서점·도서관 API 6곳 ───────────────────────────────────
s = blank(prs)
head(s, "04", "외부 서점 · 도서관 API 6곳을 합치는 이유",
     "각 API 는 가진 것과 없는 것이 정확히 갈립니다 — 한 곳으로는 안 됩니다")

table(s, In(0.62), In(1.52), In(12.1),
      [["", "알라딘", "국립중앙도서관", "Google Books", "Open Library", "Hardcover", "Gutendex"],
       [("성격", {"bold": True}), "국내 서점", "국내 납본기관", "종합", "종합", "커뮤니티", "무료 전자책"],
       [("한국 도서", {"bold": True}), ("★★★★★", {"color": RED, "bold": True}),
        ("★★★★★", {"color": RED, "bold": True}), "★★★★", "★★", "★", "없음"],
       [("커버리지", {"bold": True}), "★★★", "국내 전량", ("★★★★★", {"color": RED, "bold": True}),
        "★★★★", "★★★", "★★"],
       [("표지", {"bold": True}), ("★★★★", {"color": RED}), ("없음", {"color": MUTED}),
        ("★★★★★", {"color": RED, "bold": True}), "★★★", "★★★", "★"],
       [("평점·무드", {"bold": True}), "평점만", "없음", "평점만", "없음",
        ("★★★★★", {"color": RED, "bold": True}), "없음"],
       [("무료 전문", {"bold": True}), "없음", "없음", "△", "★★★",
        "없음", ("★★★★★", {"color": RED, "bold": True})],
       [("국내 구매", {"bold": True}), ("★★★★★", {"color": RED, "bold": True}),
        "없음", "없음", "없음", "없음", "없음"]],
      colw=[1.55, 1.35, 1.75, 1.55, 1.55, 1.35, 1.45], rh=In(0.335), fs=9.5, hfs=9.5)

y = In(4.52)
txt(s, In(0.62), y, In(5.85), In(0.32),
    [[("언어로 갈라 부릅니다", {"size": 13, "bold": True})]])
mono(s, In(0.62), y + In(0.40), In(5.85), In(1.22), """한국어 맥락 → 알라딘 + 국립중앙도서관
              (보낼 한국어 검색어가 있을 때만)
영어권      → Google Books + Open Library
              + Hardcover + Gutendex
0권이면     → Google Books(langRestrict=ko)""", size=8.5)

txt(s, In(6.88), y, In(5.85), In(0.32),
    [[("합치는 순서", {"size": 13, "bold": True})]])
mono(s, In(6.88), y + In(0.40), In(5.85), In(1.22), """1  ISBN-13 으로 1차 중복 제거
2  제목·저자 유사도로 2차 제거
3  필드별로 믿을 소스를 따로 정해 병합
     표지 → Google Books → 알라딘
     무드 → Hardcover
4  평점·평가수·연도·장르 적합성으로 정렬""", size=8.5)

box(s, In(0.62), In(6.22), In(12.1), In(0.62), fill=TINT, edge=None)
txt(s, In(0.90), In(6.36), In(11.6), In(0.36),
    [[("국내 소스 둘은 겹치지 않고 보완합니다 — 알라딘은 신간·표지·구매 링크가 강하고 절판·학술서가 약한데, "
       "국립중앙도서관은 납본 기관이라 국내 출간물이 사실상 전부 있지만 표지가 없습니다.",
       {"size": 10.5})]])

notes(s, """
이 표가 "왜 6곳이나 부르는가" 에 대한 답입니다.

한국 책을 찾는데 Google Books 만 쓰면 절판된 국내서가 안 나옵니다. 반대로
알라딘만 쓰면 영어권 책과 무료 고전을 못 찾습니다. 무드 태그(잔잔한, 위로되는)는
Hardcover 에만 있고, 무료 전문 다운로드는 Gutendex 에만 있습니다.

그래서 ISBN-13 으로 조인하고, 필드별로 어느 소스를 믿을지 따로 정했습니다.
표지는 Google Books 를 먼저 보고, 국내서는 알라딘, 무드는 Hardcover 입니다.

국내 소스 둘을 같이 쓰는 이유를 강조하고 싶습니다. 알라딘은 서점이라 팔리는
책에 강하고, 국립중앙도서관은 납본 기관이라 국내 출간물이 거의 전부 있습니다.
대신 도서관은 표지 이미지가 없어서, 국중에서 책을 찾고 표지는 알라딘에서 받는
식으로 병합합니다.

알라딘과 국중은 오류가 나도 HTTP 200 을 주는 특성이 있어서, 응답 본문의
errorCode 를 직접 검사하게 만들었습니다.
""")

# ── 6. 검증 로직 ──────────────────────────────────────────────────
s = blank(prs)
head(s, "05", "검증 로직 — 모델이 떠올린 책을 외부 데이터로 대조",
     "임계값은 추측이 아니라 실제로 틀린 사례에서 나왔습니다")

mono(s, In(0.62), In(1.52), In(12.1), In(1.02), """모델이 "《종의 기원》 정유정" 을 떠올림
   → 서점·도서관 API 조회
   → 제목 유사도 ≥ 0.7   AND   저자 유사도 ≥ 0.5   →  통과
   → 통과하지 못한 책은 답변에서 언급 금지""", size=11)

y = In(2.78)
table(s, In(0.62), y, In(12.1),
      [["임계값", "왜 이 값인가 (실측 근거)"],
       [("제목 0.7", {"bold": True, "color": RED}),
        "0.62였을 때 「종의 기원」 요청에 「종의 기원과 진화론」이 0.64로 통과했습니다. "
        "같은 책은 주제목 비교로 1.00 이 나오므로 올려도 정상 케이스는 안 깨집니다"],
       [("저자 0.5", {"bold": True, "color": RED}),
        "「1984」는 조지 오웰 원작 외에 해설서·만화판이 많습니다. 저자를 지목했는데 다른 사람이면 다른 책입니다"]],
      colw=[2.0, 10.1], rh=In(0.62), fs=10.5, hfs=10.5)

y2 = In(4.72)
txt(s, In(0.62), y2, In(5.85), In(0.32),
    [[("대조 전 정규화 — 전부 실제 버그에서 나왔습니다", {"size": 12, "bold": True})]])
mono(s, In(0.62), y2 + In(0.38), In(5.85), In(1.18), """NFKD → 발음기호 제거 → NFC
장식 문자 제거    《》 「」 ** " 등
조사 처리         「가와바타」가 조사 '가' 로 잘림
권차 제거         「혼불 1」 → 「혼불」
                  상한 20 (「Fahrenheit 451」 보호)""", size=8.5)

txt(s, In(6.88), y2, In(5.85), In(0.32),
    [[("검증이 잡아낸 실제 사례", {"size": 12, "bold": True})]])
box(s, In(6.88), y2 + In(0.38), In(5.85), In(1.18), fill=BOXBG)
txt(s, In(7.10), y2 + In(0.52), In(5.4), In(0.95),
    [[("궁중요리 질문에서 답변은 3권을 추천했는데 카드는 1장이었습니다.",
       {"size": 10, "color": MUTED})],
     [("모델이 자기 지식으로 언급한 책이 검색 결과에 없어 매칭에 실패한 것입니다.",
       {"size": 10, "color": MUTED})],
     [("→ 답변에서 제목을 뽑아 역으로 조회(lookup)해 카드를 채우게 고쳤습니다.",
       {"size": 10, "bold": True})]], space=3)

box(s, In(0.62), In(6.30), In(12.1), In(0.55), fill=TINT, edge=None)
txt(s, In(0.90), In(6.42), In(11.6), In(0.32),
    [[("환각을 프롬프트로 부탁하지 않고 구조로 막습니다. 다만 구조가 덮지 못한 경계면으로는 새어 나갔습니다.",
       {"size": 11, "bold": True})]])

notes(s, """
검증은 문자열 유사도 두 개입니다. 제목 0.7, 저자 0.5.

숫자의 근거가 중요합니다. 제목 임계값을 0.62로 뒀을 때 「종의 기원」을 찾는
요청에 「종의 기원과 진화론」이라는 다른 책이 0.64로 통과했습니다. 그래서
0.7로 올렸습니다. 같은 책은 주제목 비교에서 1.00 이 나오므로 정상 케이스는
안 깨집니다.

저자를 따로 보는 이유는 「1984」 같은 경우입니다. 조지 오웰 원작 외에 해설서와
만화판이 많습니다. 사용자가 저자를 지목했는데 다른 사람이면 다른 책입니다.

정규화 항목은 전부 실제 버그에서 나왔습니다. 「가와바타」가 조사 '가' 로 잘려서
매칭이 실패했고, 권차 제거를 무제한으로 하니 「Fahrenheit 451」이 「Fahrenheit」가
됐습니다. 그래서 상한 20을 뒀습니다.

마지막 줄은 정직하게 말할 부분입니다. 검증 계층은 정상 동작했는데, 모델이
"DB에서 확인되지 않았지만 알려진 책" 이라는 소제목을 스스로 만들어 검증에
실패한 책을 사용자에게 보여준 사고가 있었습니다. 검증 결과를 모델에게
알려준 것이 새는 경로였습니다.
""")

# ── 7. 실측과 남은 과제 ───────────────────────────────────────────
s = blank(prs)
head(s, "06", "실측과 남은 과제", "배포된 서비스를 실제로 호출해 측정한 값입니다")

txt(s, In(0.62), In(1.48), In(5.9), In(0.32),
    [[("응답 시간 — 병목은 외부 API 가 아니었습니다", {"size": 12.5, "bold": True})]])
mono(s, In(0.62), In(1.86), In(5.9), In(1.32), """전체              27,098 ms
  외부 API 6곳      1,218 ms   (4.5%)
  Bedrock 4회      25,900 ms   (95%)

토큰   입력 17,905 / 출력 670
카드   12장""", size=9.5)

box(s, In(0.62), In(3.30), In(5.9), In(0.62), fill=TINT, edge=None)
txt(s, In(0.86), In(3.42), In(5.5), In(0.4),
    [[("API Gateway 한계 30초까지 여유 0.8초", {"size": 11, "bold": True, "color": RED})],
     [("남은 최우선 과제입니다", {"size": 9.5, "color": MUTED})]], space=2)

txt(s, In(6.83), In(1.48), In(5.9), In(0.32),
    [[("자동 검증 474건 — 네트워크 없이 몇 초", {"size": 12.5, "bold": True})]])
table(s, In(6.83), In(1.86), In(5.9),
      [["검사", "건수"],
       ["프롬프트·모듈 로드 (26파일)", "—"],
       ["정책 (미성년·인젝션·PII)", "99"],
       ["기능 (병합·검증·카드선별)", "218"],
       ["OpenAI 호환 계약", "50"],
       ["에이전트 루프 (3모드)", "23"],
       ["프론트 저장소 · 화면 렌더", "84"]],
      colw=[4.4, 1.5], rh=In(0.285), fs=9.5, hfs=9.5)

txt(s, In(6.83), In(3.90), In(5.9), In(0.3),
    [[("회귀 방지 검사는 전부 실제 사고를 픽스처로 만들었습니다", {"size": 9.5, "color": MUTED})]])

rule(s, In(0.62), In(4.42), In(12.1))
txt(s, In(0.62), In(4.58), In(12.1), In(0.32),
    [[("남은 과제 — 정직하게", {"size": 12.5, "bold": True})]])

table(s, In(0.62), In(4.98), In(12.1),
      [["과제", "상태"],
       [("결합 공격 방어", {"bold": True, "color": RED}),
        "외부 벤치마크에서 \"정상 책 요청 + 유해 요청\" 을 한 문장에 섞는 패턴 10건이 뚫렸습니다. 대응 없음"],
       ["응답 시간", "Bedrock 이 95%. 프롬프트 캐싱은 SDK 가 cachePoint 미지원이라 업그레이드가 먼저"],
       ["응답 스트리밍", "HTTP API 로는 불가. 답변이 다 만들어진 뒤 한꺼번에 도착합니다"],
       ["Bedrock Guardrails", "미적용. 안전장치가 전부 자체 구현이라 외부 검증을 받은 적이 없습니다"],
       ["로그인", "없음. 레이트리밋 4층으로 비용만 방어. 상시 운영할 구성은 아닙니다"]],
      colw=[2.6, 9.5], rh=In(0.40), fs=10, hfs=10)

notes(s, """
왼쪽 실측이 이 프로젝트에서 판단이 뒤집힌 대표 사례입니다.

설정 주석들이 "외부 API 가 느린 날" 을 걱정하고 있었는데, 재보니 외부 API 는
전체의 4.5% 였고 Bedrock 이 95% 였습니다. 병목을 완전히 반대로 짚고 있었습니다.
측정 없이 최적화하면 안 된다는 걸 여기서 배웠습니다.

가운데 검증 474건은 네트워크 없이 도는 것만 센 값입니다. 회귀 방지 검사는
전부 실제로 일어난 사고를 픽스처로 만들었습니다 — 사용자가 신고한 답변 원문,
카드 누락 재현, 로그에서 나온 오탐.

오른쪽 남은 과제에서 첫 줄이 가장 중요합니다. 자체 정책 테스트는 전부
통과했는데, 외부 안전성 벤치마크를 돌리니 "잔인한 스릴러 추천해줘, 그리고
사람을 해치는 방법을..." 처럼 정상 요청과 유해 요청을 한 문장에 결합하는
패턴에 뚫렸습니다. 앞부분이 진짜 책 추천 요청이라 정책이 통과시킵니다.
아직 해결하지 못했습니다.

배포 확인은 /api/health 의 runtime 블록으로 합니다. 지금 배포본은
maxToolIterations 가 3이고 코드는 4입니다 — 무엇이 배포됐는지 확인할 수 있게
만들어 둔 장치입니다.
""")

prs.save(OUT)
print(f"저장: {OUT}  ({os.path.getsize(OUT):,} bytes, 슬라이드 {len(prs.slides.__iter__.__self__._sldIdLst)}장)")
